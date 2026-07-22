"""
AI Boardroom — Extractor
Extracts raw text from various document formats (PDF, DOCX, PPTX).
"""

from __future__ import annotations

from io import BytesIO

import docx
import csv
import openpyxl
import PyPDF2
from pptx import Presentation

from app.core.logging import get_logger

logger = get_logger(__name__)


def extract_text(file_content: bytes, mime_type: str, filename: str) -> str:
    """Route extraction based on MIME type."""
    try:
        if mime_type == "application/pdf":
            return _extract_pdf(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_docx(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_pptx(file_content)
        elif mime_type == "text/csv" or filename.endswith(".csv"):
            return _extract_csv(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _extract_xlsx(file_content)
        elif mime_type == "application/json":
            return file_content.decode("utf-8")
        elif mime_type == "text/plain":
            return file_content.decode("utf-8")
        else:
            return ""
    except Exception as e:
        logger.error("Failed to extract text", filename=filename, error=str(e))
        return ""


def _extract_pdf(content: bytes) -> str:
    reader = PyPDF2.PdfReader(BytesIO(content))
    text = []
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text.append(extracted)
    return "\n\n".join(text)


def _extract_docx(content: bytes) -> str:
    doc = docx.Document(BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs])


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n\n".join(text)


def _extract_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    reader = csv.reader(text.splitlines())
    rows = []
    for row in reader:
        rows.append("\t".join(row))
    return "\n".join(rows)


def _extract_xlsx(content: bytes) -> str:
    wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            if any(row_text):
                text.append("\t".join(row_text))
    return "\n".join(text)
